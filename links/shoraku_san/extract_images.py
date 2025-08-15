from pptx import Presentation
import os

# PPTX file ka path
pptx_path = "ggggg.pptx"
output_dir = "images"

# Images ke liye folder banao
os.makedirs(output_dir, exist_ok=True)

# PPTX file kholo
prs = Presentation(pptx_path)

# Har slide se images extract karo
image_count = 1
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            image = shape.image
            image_bytes = image.blob
            # Image ka naam set karo (image1.png, image2.png, etc.)
            image_filename = f"{output_dir}/image{image_count}.png"
            with open(image_filename, "wb") as f:
                f.write(image_bytes)
            image_count += 1
print(f"Extracted {image_count-1} images to {output_dir} folder")